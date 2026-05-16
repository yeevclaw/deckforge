"""
Convert a directory of HTML slide pages to a PPTX deck.

Each .html file in --pages-dir becomes one slide (16:9, 1280×720 logical size).
Pages are ordered by filename (so name them page_01.html, page_02.html, ...).

If a planning.json sits next to the pages dir (or is passed via --planning),
speaker_notes are attached to each slide.

Modes:
  --mode image     (default) — render each page to a high-DPI PNG, embed as full-bleed slide background.
                                Perfect visual fidelity. Text is not editable in PowerPoint.
  --mode editable  — basic shape reconstruction (placeholder; image mode recommended)

Examples:
  python html_to_pptx.py --pages-dir pages/ --output deck.pptx
  python html_to_pptx.py --pages-dir pages/ --output deck.pptx --planning planning.json
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
except ImportError:
    sys.stderr.write("python-pptx not installed. Run: pip install python-pptx --break-system-packages\n")
    sys.exit(1)

# Local import
SCRIPT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPT_DIR))
from render_html import render  # noqa: E402

# 16:9 standard slide @ 96 dpi == 1280×720
SLIDE_W_EMU = Emu(int(13.333 * 914400))   # 13.333 inches
SLIDE_H_EMU = Emu(int(7.5 * 914400))      # 7.5 inches


def collect_pages(pages_dir: Path) -> list[Path]:
    pages = sorted(p for p in pages_dir.iterdir() if p.suffix.lower() == ".html")
    if not pages:
        raise FileNotFoundError(f"No .html files found in {pages_dir}")
    return pages


def load_planning(planning_path: Path | None, pages_dir: Path) -> dict:
    if planning_path and planning_path.exists():
        return json.loads(planning_path.read_text(encoding="utf-8"))
    # Default location
    default = pages_dir.parent / "planning.json"
    if default.exists():
        return json.loads(default.read_text(encoding="utf-8"))
    return {}


def speaker_notes_for(page_idx: int, planning: dict) -> str:
    pages = planning.get("pages") or []
    if 0 <= page_idx < len(pages):
        return pages[page_idx].get("speaker_notes", "") or ""
    return ""


def build_image_mode(pages: list[Path], out_path: Path, planning: dict,
                     workdir: Path, scale: float = 2.0, engine: str = "auto"):
    workdir.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    blank_layout = prs.slide_layouts[6]  # blank

    for i, html_path in enumerate(pages):
        png_path = workdir / f"{html_path.stem}.png"
        print(f"[{i+1}/{len(pages)}] Rendering {html_path.name} → PNG ...", flush=True)
        engine_used = render(html_path, png_path, 1280, 720, scale, engine)
        if i == 0:
            print(f"  (engine: {engine_used})")

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(png_path), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)

        notes = speaker_notes_for(i, planning)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(out_path))
    print(f"Wrote {out_path}")


def build_editable_mode(pages, out_path, planning, workdir, scale, engine):
    """
    Stub for editable mode. Currently a hybrid: it embeds the rendered image as
    background AND adds a transparent text-box overlay for the page title region
    (top 12% of the slide), seeded from planning.json if available.

    Full HTML→shape reconstruction is a longer-term project. For now, this lets
    users edit titles after handoff while keeping pixel-perfect bodies.
    """
    workdir.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU
    blank_layout = prs.slide_layouts[6]

    planning_pages = planning.get("pages") or []

    for i, html_path in enumerate(pages):
        png_path = workdir / f"{html_path.stem}.png"
        print(f"[{i+1}/{len(pages)}] Rendering {html_path.name} → PNG (editable mode) ...", flush=True)
        render(html_path, png_path, 1280, 720, scale, engine)

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(png_path), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)

        # Title overlay: top-left, width 80%, height 12% of slide.
        # Transparent background, large bold text.
        if i < len(planning_pages):
            title = planning_pages[i].get("title", "")
            if title:
                from pptx.dml.color import RGBColor
                left = Inches(0.5)
                top = Inches(0.3)
                width = Inches(10)
                height = Inches(0.9)
                tb = slide.shapes.add_textbox(left, top, width, height)
                tf = tb.text_frame
                tf.text = title
                p = tf.paragraphs[0]
                p.font.size = Pt(32)
                p.font.bold = True
                # Don't set color — let the user override. Pixel beneath still shows.
                # Make text invisible by default to avoid double-rendering; user reveals when editing.
                p.font.color.rgb = RGBColor(0, 0, 0)
                # Hack: make the text "transparent" by setting fill type to no fill
                # (Actually with python-pptx the simplest is to set alpha — not supported.
                # So we leave it visible and let the renderer pick — the rendered title
                # underneath is at top 8% so we offset our overlay slightly to avoid overlap.)

        notes = speaker_notes_for(i, planning)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(out_path))
    print(f"Wrote {out_path}")
    print("NOTE: editable mode places title overlays on top. You may want to delete the "
          "rendered title in the background image if it conflicts. For full editability, "
          "the recommended workflow is to edit titles in the HTML and re-render.")


def main():
    p = argparse.ArgumentParser(description=__doc__,
                                formatter_class=argparse.RawDescriptionHelpFormatter)
    p.add_argument("--pages-dir", required=True, help="Directory containing page_*.html files")
    p.add_argument("--output", required=True, help="Output .pptx path")
    p.add_argument("--planning", default=None, help="planning.json path (auto-detected if omitted)")
    p.add_argument("--mode", default="image", choices=["image", "editable"])
    p.add_argument("--scale", type=float, default=2.0,
                   help="Rendering scale factor (2.0 → 2560×1440 PNG)")
    p.add_argument("--engine", default="auto",
                   choices=["auto", "playwright", "libreoffice"],
                   help="HTML rendering engine")
    p.add_argument("--workdir", default=None,
                   help="Directory for intermediate PNGs (default: <pages-dir>/_renders)")
    args = p.parse_args()

    pages_dir = Path(args.pages_dir).resolve()
    out_path = Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    planning_path = Path(args.planning).resolve() if args.planning else None
    workdir = Path(args.workdir).resolve() if args.workdir else pages_dir / "_renders"

    pages = collect_pages(pages_dir)
    planning = load_planning(planning_path, pages_dir)

    if args.mode == "image":
        build_image_mode(pages, out_path, planning, workdir, args.scale, args.engine)
    else:
        build_editable_mode(pages, out_path, planning, workdir, args.scale, args.engine)


if __name__ == "__main__":
    main()
