"""
Convert a directory of SVG slide pages to a PPTX deck.

Each .svg file in --pages-dir becomes one slide (16:9, 1280×720 viewBox).
Pages are ordered by filename (so name them page_01.svg, page_02.svg, ...).

If a planning.json sits next to the pages dir (or is passed via --planning),
speaker_notes are attached to each slide.

By default each slide is split into two stacked pictures:
  1. A full rasterized PNG render — the movable "background image" that every
     viewer shows correctly (Keynote, macOS Preview, Quick Look,
     PowerPoint <2016, web-mail previewers, etc.). It carries the atmosphere
     (gradient backgrounds, glassmorphism, soft shadows) that PowerPoint can't
     turn into editable shapes anyway.
  2. A transparent CONTENT-ONLY layer carrying the original geometry via the
     PowerPoint 2016+ svgBlip OOXML extension — atmosphere removed, card shadow
     <filter>s stripped, and <use>/<symbol> icons inlined into real geometry.
     Right-click → "Convert to Shape" on this layer makes the text, cards,
     icons and lines individually editable / movable / resizable.
Pass --no-decompose to embed each slide as a single picture (full render +
svgBlip) the way older versions did.

The PNG fallback is rendered using whichever SVG renderer is available,
in this preference order:
  resvg-py → cairosvg → inkscape → rsvg-convert

resvg-py is the only one with **zero system dependencies** — it ships a
pure Rust binary inside a pip wheel, available for Python 3.9+ on
macOS / Linux / Windows. That's why setup.sh installs it by default.

If none of these is installed, a 1×1 transparent placeholder PNG is used
(sufficient for PowerPoint 2016+ but will look blank in Keynote /
macOS Preview) and a warning is printed.

Required dependency:    python-pptx       (pip install python-pptx)
Strongly recommended:   resvg-py          (pip install resvg-py)
                        img2pdf           (pip install img2pdf)
                                          enables the companion .pdf alongside the .pptx
Alternatives:           cairosvg          pip install cairosvg + libcairo
                        rsvg-convert      brew install librsvg / apt librsvg2-bin
                        inkscape          brew install inkscape

Flags:
  --no-svg            Skip the svgBlip extension; PPTX becomes image-only.
                      Use this if a viewer chokes on the SVG ext (some
                      older Keynote versions).
  --placeholder-only  Force the 1×1 transparent placeholder PNG even when
                      a real renderer is available. Smaller file, but the
                      slide only displays in PowerPoint 2016+.
  --png-width N       Width (px) of the PNG fallback. Default 2560 (2× DPI).
  --no-anim           Disable flow-edge animation. By default, a page whose SVG
                      marks dashed edges with class="flow-anim" is embedded as a
                      looping animated GIF (dashes flow in PowerPoint / Keynote
                      slideshow mode). Animated slides skip the svgBlip layer,
                      so they are not Convert-to-Shape editable.
  --gif-width N       Width (px) of animated GIF frames. Default 1600.
  --no-decompose      Embed each slide as one picture (full render + svgBlip)
                      instead of the default background-image + editable-content
                      split.

Examples:
  python svg_to_pptx.py --pages-dir pages/ --output deck.pptx
  python svg_to_pptx.py --pages-dir pages/ --output deck.pptx --no-svg
  python svg_to_pptx.py --pages-dir pages/ --output deck.pptx --planning planning.json
"""
from __future__ import annotations

import argparse
import base64
import json
import re
import shutil
import subprocess
import sys
from copy import deepcopy
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError as e:
    sys.stderr.write(
        f"Missing dependency: {e}\n"
        "Install: pip install python-pptx --break-system-packages\n"
        "(python-pptx pulls in lxml + Pillow automatically.)\n"
    )
    sys.exit(1)

# img2pdf is optional — used to also emit a .pdf alongside the .pptx.
try:
    import img2pdf  # type: ignore
    _HAS_IMG2PDF = True
except ImportError:
    _HAS_IMG2PDF = False

SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
SVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# SVG document namespaces (the slide files themselves, not the OOXML wrappers).
SVG_DOC_NS = "http://www.w3.org/2000/svg"
XLINK_NS = "http://www.w3.org/1999/xlink"

# 16:9 standard slide @ 96 dpi == 1280×720 logical
SLIDE_W_EMU = Emu(int(13.333 * 914400))
SLIDE_H_EMU = Emu(int(7.5 * 914400))

# A valid 1×1 fully-transparent PNG. PowerPoint stretches it to fill the slide;
# since it's transparent, the SVG vector layer above is what the viewer sees.
# This satisfies the OOXML requirement that <a:blip> reference an image.
_PLACEHOLDER_PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlE"
    "QVR42mNk+P+/HgAFhAJ/wlseKgAAAABJRU5ErkJggg=="
)


# ---------- SVG → PNG rendering ----------

def _render_with_resvg_py(svg_path: Path, png_path: Path, width: int) -> bool:
    try:
        import resvg_py
    except ImportError:
        return False
    svg_data = svg_path.read_text(encoding="utf-8")
    try:
        png_bytes = resvg_py.svg_to_bytes(svg_string=svg_data, width=width)
    except Exception as e:
        print(f"[svg_to_pptx] resvg-py failed on {svg_path.name}: {e}", file=sys.stderr)
        return False
    png_path.write_bytes(bytes(png_bytes))
    return True


def _render_with_cairosvg(svg_path: Path, png_path: Path, width: int) -> bool:
    try:
        import cairosvg
    except ImportError:
        return False
    try:
        cairosvg.svg2png(url=str(svg_path), write_to=str(png_path), output_width=width)
    except Exception as e:
        # cairosvg installs cleanly but fails at runtime if libcairo C lib is missing.
        print(f"[svg_to_pptx] cairosvg failed on {svg_path.name}: {e}", file=sys.stderr)
        return False
    return True


def _render_with_inkscape(svg_path: Path, png_path: Path, width: int) -> bool:
    ink = shutil.which("inkscape")
    if not ink:
        return False
    try:
        subprocess.run(
            [ink, "--export-type=png", f"--export-filename={png_path}",
             f"--export-width={width}", str(svg_path)],
            check=True, capture_output=True, timeout=30,
        )
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as e:
        print(f"[svg_to_pptx] inkscape failed: {e}", file=sys.stderr)
        return False
    return png_path.exists()


def _render_with_rsvg(svg_path: Path, png_path: Path, width: int) -> bool:
    rsvg = shutil.which("rsvg-convert")
    if not rsvg:
        return False
    try:
        subprocess.run(
            [rsvg, "-w", str(width), "-o", str(png_path), str(svg_path)],
            check=True, capture_output=True, timeout=30,
        )
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as e:
        print(f"[svg_to_pptx] rsvg-convert failed: {e}", file=sys.stderr)
        return False
    return png_path.exists()


def render_real_png(svg_path: Path, png_path: Path, width: int = 2560):
    """Render a high-DPI PNG fallback. Returns engine name, or None if no renderer."""
    png_path.parent.mkdir(parents=True, exist_ok=True)
    for name, fn in (
        ("resvg-py", _render_with_resvg_py),     # pure-pip, no system deps
        ("cairosvg", _render_with_cairosvg),     # needs libcairo
        ("inkscape", _render_with_inkscape),     # separate binary
        ("rsvg-convert", _render_with_rsvg),     # native lib
    ):
        if fn(svg_path, png_path, width):
            return name
    return None


def verify_renderer_works(test_svg: Path) -> str | None:
    """Do one real render to confirm at least one renderer actually works
    (catches cairosvg-installed-but-libcairo-missing-style failures up front)."""
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        tmp_path = Path(tmp.name)
    try:
        return render_real_png(test_svg, tmp_path, width=64)
    finally:
        tmp_path.unlink(missing_ok=True)


# ---------- Flow-edge animated GIF rendering ----------
#
# A page opts in by marking dashed edges with class="flow-anim" (must also
# carry a stroke-dasharray). The page is then rendered once per frame with a
# shifted stroke-dashoffset and assembled into a seamlessly looping GIF —
# total offset travel per cycle equals the dasharray period, so frame N wraps
# exactly back to frame 0. PowerPoint and Keynote play GIFs in slideshow mode
# (embedded SVG animations do NOT play — svgBlip renders statically), which is
# why animated slides embed a GIF and skip the svgBlip layer.

FLOW_ANIM_CLASS = "flow-anim"
GIF_FRAMES = 12
GIF_FRAME_MS = 80   # GIF timing is centisecond-granular; keep it a multiple of 10
GIF_TRAVEL_PERIODS = 2  # dash travel per loop, in periods — any integer keeps the
                        # loop seamless; 2 ≈ 29px/s for "8 6" (1 looks sleepy,
                        # 3 drops to 4 samples/period and stutters at 12 frames)
DEFAULT_GIF_WIDTH = 1600


def _svg_has_flow_anim_class(svg_text: str) -> bool:
    """True if any element carries class="flow-anim" (regardless of dasharray)."""
    return re.search(
        r'class\s*=\s*["\'][^"\']*\b' + FLOW_ANIM_CLASS + r'\b', svg_text
    ) is not None


def _flow_anim_is_animatable(svg_text: str, name: str) -> bool:
    """A page becomes a GIF only if some element carries BOTH class="flow-anim"
    AND a stroke-dasharray — the dashes that actually flow. A flow-anim class with
    no dasharray animates nothing yet still costs the slide its editable layer (GIF
    slides skip svgBlip), so we refuse GIF mode, keep the page static and
    Convert-to-Shape editable, and warn loudly. This enforces the contract in
    prompts/05_designer_svg.md Step 5.7 ("flow-anim + stroke-dasharray") and stops a
    stray or mislabeled class from silently turning a whole page into a non-editable
    GIF."""
    if not _svg_has_flow_anim_class(svg_text):
        return False
    for tag in re.finditer(r"<[^>]*>", svg_text):
        t = tag.group(0)
        if re.search(r'class\s*=\s*["\'][^"\']*\b' + FLOW_ANIM_CLASS + r'\b', t) and \
                re.search(r'stroke-dasharray\s*=\s*["\'][^"\']+["\']', t):
            return True
    print(f"[svg_to_pptx] ⚠️  {name}: class=\"flow-anim\" is present but no element "
          "also carries a stroke-dasharray — nothing would animate. Keeping the slide "
          "STATIC and Convert-to-Shape editable (no GIF). Add a stroke-dasharray to the "
          "open <line>/<path> to animate it, or remove the flow-anim class.",
          file=sys.stderr)
    return False


def _flow_dash_period(svg_text: str) -> float:
    """Sum of the first flow-anim element's stroke-dasharray values — the px
    distance one dash cycle travels. Falls back to 14.0 (dasharray "8 6") on
    any parse failure; a wrong period only makes the loop seam visible, it
    never breaks rendering."""
    for pattern in (
        # class before stroke-dasharray, and the reverse attribute order
        r'<[^>]*\b' + FLOW_ANIM_CLASS + r'\b[^>]*stroke-dasharray\s*=\s*["\']([^"\']+)["\']',
        r'<[^>]*stroke-dasharray\s*=\s*["\']([^"\']+)["\'][^>]*\b' + FLOW_ANIM_CLASS + r'\b',
    ):
        m = re.search(pattern, svg_text)
        if m:
            try:
                total = sum(float(v) for v in re.split(r"[\s,]+", m.group(1).strip()) if v)
                if total > 0:
                    return total
            except ValueError:
                pass
    return 14.0


def _inject_dashoffset_style(svg_text: str, offset_px: float) -> str:
    """Insert a <style> rule right after the opening <svg> tag. CSS author
    styles override presentation attributes, so this wins over any
    stroke-dashoffset already on the elements (verified with resvg)."""
    style = f"<style>.{FLOW_ANIM_CLASS}{{stroke-dashoffset:{offset_px:.3f}px;}}</style>"
    return re.sub(r"(<svg\b[^>]*>)", r"\1" + style, svg_text, count=1)


def _lint_flow_anim(svg_text: str, period: float, name: str) -> None:
    """Non-fatal design checks on a flow-anim page (warnings to stderr).
    The GIF still renders either way — these catch the two mistakes that
    produce a working-but-ugly animation."""
    tags = [m.group(0) for m in re.finditer(r"<[^>]*>", svg_text)
            if re.search(r'class\s*=\s*["\'][^"\']*\b' + FLOW_ANIM_CLASS + r"\b",
                         m.group(0))]
    # One dashoffset is injected for the whole class, so mixed dasharrays loop
    # with a visible seam on every element except the first.
    dasharrays = set()
    for tag in tags:
        m = re.search(r'stroke-dasharray\s*=\s*["\']([^"\']+)["\']', tag)
        if m:
            dasharrays.add(" ".join(m.group(1).replace(",", " ").split()))
    if len(dasharrays) > 1:
        print(f"[svg_to_pptx] ⚠️  {name}: flow-anim elements mix stroke-dasharray "
              f"values {sorted(dasharrays)} — only the first defines the loop "
              "period, the others will show a seam each cycle. Use ONE dasharray "
              "per page.", file=sys.stderr)
    # Closed shapes animate into marching-ants selection boxes (diagrams.md ban).
    # An orbit / cycle ring must be built from open <path> arcs, never one closed
    # <circle>/<rect>/<polygon>, so the dashes read as rotation, not a marquee.
    for tag in tags:
        m = re.match(r"<\s*([A-Za-z]+)", tag)
        if m and m.group(1).lower() in ("circle", "ellipse", "rect", "polygon"):
            print(f"[svg_to_pptx] ⚠️  {name}: flow-anim on a <{m.group(1)}> (a closed "
                  "shape) animates into a marching-ants selection box. Use open "
                  "<line>/<path> arcs instead (an orbit ring = several open arcs). "
                  "See references/diagrams.md → 'never animate'.", file=sys.stderr)
    # Flowing dashes need room to read as flow: ≥5 periods of visible length.
    min_len = 5 * period
    for tag in tags:
        if not tag.startswith("<line"):
            continue  # path lengths aren't worth computing here
        try:
            x1, y1, x2, y2 = (float(re.search(c + r'\s*=\s*["\']([^"\']+)["\']', tag).group(1))
                              for c in ("x1", "y1", "x2", "y2"))
        except AttributeError:
            continue
        length = ((x2 - x1) ** 2 + (y2 - y1) ** 2) ** 0.5
        if length < min_len:
            print(f"[svg_to_pptx] ⚠️  {name}: flow-anim line of {length:.0f}px is "
                  f"shorter than 5 dash periods ({min_len:.0f}px) — too short to "
                  "read as flow; consider a static connector instead.",
                  file=sys.stderr)


def render_flow_anim_gif(svg_path: Path, gif_path: Path, width: int,
                         frames: int = GIF_FRAMES,
                         frame_ms: int = GIF_FRAME_MS) -> bool:
    """Render svg_path into a looping GIF at gif_path. Requires resvg-py
    (the only renderer with in-memory svg_string input). Returns False on any
    failure so the caller falls back to the normal static slide."""
    try:
        import resvg_py
        from io import BytesIO
        from PIL import Image
    except ImportError as e:
        print(f"[svg_to_pptx] flow-anim GIF needs resvg-py + Pillow ({e}); "
              f"{svg_path.name} stays static.", file=sys.stderr)
        return False
    try:
        svg_text = svg_path.read_text(encoding="utf-8")
        period = _flow_dash_period(svg_text)
        _lint_flow_anim(svg_text, period, svg_path.name)
        travel = period * GIF_TRAVEL_PERIODS
        rgb_frames = []
        for i in range(frames):
            offset = -(travel * i / frames)  # negative → dashes flow in path direction
            png_bytes = resvg_py.svg_to_bytes(
                svg_string=_inject_dashoffset_style(svg_text, offset), width=width)
            rgb_frames.append(Image.open(BytesIO(bytes(png_bytes))).convert("RGB"))
        # Quantize all frames against one shared palette, every frame with
        # dither=NONE (including frame 0 — quantize() defaults to
        # Floyd-Steinberg, and a dithered first frame flickers once per loop).
        pal = rgb_frames[0].quantize(colors=256, method=Image.Quantize.MEDIANCUT)
        quantized = [f.quantize(palette=pal, dither=Image.Dither.NONE)
                     for f in rgb_frames]
        gif_path.parent.mkdir(parents=True, exist_ok=True)
        quantized[0].save(str(gif_path), format="GIF", save_all=True,
                          append_images=quantized[1:], duration=frame_ms, loop=0,
                          disposal=1, optimize=True)
        return True
    except Exception as e:
        print(f"[svg_to_pptx] flow-anim GIF failed on {svg_path.name}: {e}; "
              "slide stays static.", file=sys.stderr)
        return False


# ---------- PPTX assembly ----------

def collect_pages(pages_dir: Path) -> list[Path]:
    pages = sorted(p for p in pages_dir.iterdir() if p.suffix.lower() == ".svg")
    if not pages:
        raise FileNotFoundError(f"No .svg files found in {pages_dir}")
    return pages


def load_planning(planning_path: Path | None, pages_dir: Path) -> dict:
    if planning_path and planning_path.exists():
        return json.loads(planning_path.read_text(encoding="utf-8"))
    default = pages_dir.parent / "planning.json"
    if default.exists():
        return json.loads(default.read_text(encoding="utf-8"))
    return {}


def speaker_notes_for(idx: int, planning: dict) -> str:
    pages = planning.get("pages") or []
    if 0 <= idx < len(pages):
        return pages[idx].get("speaker_notes", "") or ""
    return ""


def _next_partname(package, base: str, ext: str) -> PackURI:
    existing = {p.partname for p in package.iter_parts()}
    n = 1
    while True:
        candidate = PackURI(f"/ppt/media/{base}{n}.{ext}")
        if candidate not in existing:
            return candidate
        n += 1


def _add_svg_part(slide_part, svg_bytes: bytes) -> str:
    package = slide_part.package
    partname = _next_partname(package, "image", "svg")
    svg_part = Part(
        partname=partname,
        content_type="image/svg+xml",
        blob=svg_bytes,
        package=package,
    )
    return slide_part.relate_to(svg_part, RT.IMAGE)


def _inject_svg_ext(pic_elem, svg_rId: str):
    """Append <a:extLst><a:ext uri=...><asvg:svgBlip r:embed=svg_rId/></a:ext></a:extLst>
    to the <a:blip> inside the picture, so Office 2016+ renders the SVG vector."""
    blip = pic_elem.find('.//' + qn('a:blip'))
    if blip is None:
        return
    existing = blip.find(qn('a:extLst'))
    if existing is not None:
        blip.remove(existing)
    extLst = etree.SubElement(blip, qn('a:extLst'))
    ext = etree.SubElement(extLst, qn('a:ext'))
    ext.set('uri', SVG_EXT_URI)
    svgBlip = etree.SubElement(ext, f'{{{SVG_NS}}}svgBlip', nsmap={'asvg': SVG_NS})
    svgBlip.set(f'{{{R_NS}}}embed', svg_rId)


# ---------- Editable-content layer (Convert-to-Shape friendly) ----------
#
# PowerPoint's "Convert to Shape" only turns a subset of SVG into native,
# independently movable shapes: plain <text>, <rect>, <circle>, <path>, <line>.
# Anything using features it can't represent — blur/shadow <filter>s, gradient
# washes, glassmorphism, and (unreliably) <use>/<symbol> icon references — is
# rasterized into the converted picture, so the user can't move or resize it.
#
# So each slide is built as TWO stacked pictures:
#   A. the full PNG render (every viewer, incl. Keynote, shows the complete
#      slide; this is the one movable "background image").
#   B. a transparent CONTENT-ONLY layer carrying the svgBlip vector — atmosphere
#      removed (it lives in A), card <filter> shadows stripped (the soft shadow
#      still shows from A; the card itself becomes a clean editable shape), and
#      <use>/<symbol> icons inlined into real geometry. Convert-to-Shape on B
#      yields editable, movable text / cards / icons / lines.
# B is layered exactly over A, so the slide looks identical; A guarantees the
# look in every viewer, so a misclassified element is at worst non-editable,
# never visually wrong.

_PAINTED_TAGS = {"rect", "circle", "ellipse", "path", "polygon", "polyline",
                 "line", "image"}
_DEF_CONTAINERS = {"defs", "symbol", "marker", "clipPath", "mask", "pattern",
                   "linearGradient", "radialGradient", "filter"}


def _ln(el) -> str:
    """Local (namespace-stripped) tag name of an lxml element. Returns "" for
    comments / processing instructions (whose .tag is not a string)."""
    if not isinstance(el.tag, str):
        return ""
    return etree.QName(el).localname


def _covers_canvas(w: str | None, h: str | None) -> bool:
    """True if a width/height pair spans (almost) the whole 1280×720 canvas."""
    def big(v, full):
        if v is None:
            return False
        if v.strip() in ("100%",):
            return True
        try:
            return float(v) >= full * 0.95
        except ValueError:
            return False
    return big(w, 1280) and big(h, 720)


def _is_atmosphere(el) -> bool:
    """Decide whether a painted element belongs to the (non-editable) background
    atmosphere — and so should be dropped from the editable layer and shown only
    via the full background render. Bias toward True: anything excluded here is
    still painted by background picture A, so the look never changes; it just
    won't be vector-editable."""
    if _ln(el) == "text":
        return False  # text is always editable content
    # Explicit marker is authoritative — on the element OR any ancestor <g>, so a
    # whole decorative group can be forced into the background image with one class.
    # See references/editable_mode.md.
    if "atmosphere" in (el.get("class") or ""):
        return True
    if any("atmosphere" in (a.get("class") or "") for a in el.iterancestors()):
        return True
    if (el.get("fill") or "").startswith("url("):
        return True  # gradient fill → background wash / glow / glass
    for attr in ("opacity", "fill-opacity"):
        v = el.get(attr)
        if v is not None:
            try:
                if float(v) < 1.0:
                    return True  # translucent → decorative layer
            except ValueError:
                pass
    if _ln(el) == "rect":
        x = (el.get("x") or "0").strip()
        y = (el.get("y") or "0").strip()
        if x in ("0", "0.0") and y in ("0", "0.0") \
                and _covers_canvas(el.get("width"), el.get("height")):
            return True  # full-canvas background plate
    return False


def _inline_uses(root, symbol_map: dict) -> None:
    """Replace every <use href="#icon-…"> with the referenced <symbol>'s inlined
    geometry, so PowerPoint converts it to movable freeform shapes instead of
    rasterizing the reference. Positioning/scaling from the <use> x/y/width/height
    and the symbol viewBox is reproduced with a transform; the <use> color is set
    on the wrapper <g> so the icons' stroke="currentColor" still resolves."""
    use_tag = f"{{{SVG_DOC_NS}}}use"
    for use in list(root.iter(use_tag)):
        href = use.get("href") or use.get(f"{{{XLINK_NS}}}href")
        if not href or not href.startswith("#"):
            continue
        sym = symbol_map.get(href[1:])
        if sym is None:
            continue
        parent = use.getparent()
        if parent is None:
            continue
        g = parent.makeelement(f"{{{SVG_DOC_NS}}}g", {})
        parts = []
        x = float(use.get("x") or 0)
        y = float(use.get("y") or 0)
        if x or y:
            parts.append(f"translate({x} {y})")
        vb = sym.get("viewBox")
        w, h = use.get("width"), use.get("height")
        if vb and w and h:
            try:
                vbx, vby, vbw, vbh = (float(t) for t in re.split(r"[ ,]+", vb.strip()))
                sx = float(w) / vbw if vbw else 1.0
                sy = float(h) / vbh if vbh else 1.0
                if sx != 1.0 or sy != 1.0:
                    parts.append(f"scale({sx} {sy})")
                if vbx or vby:
                    parts.append(f"translate({-vbx} {-vby})")
            except (ValueError, ZeroDivisionError):
                pass
        if parts:
            g.set("transform", " ".join(parts))
        if use.get("color"):
            g.set("color", use.get("color"))
        for child in sym:
            g.append(deepcopy(child))
        use.addprevious(g)
        parent.remove(use)


def build_editable_layer(svg_bytes: bytes) -> bytes | None:
    """Return a content-only SVG (atmosphere removed, filters stripped, icons
    inlined) for the svgBlip editable layer. Returns None if parsing fails or
    nothing usable remains — the caller then falls back to the single-picture
    slide so behavior never regresses."""
    try:
        root = etree.fromstring(svg_bytes)
    except etree.XMLSyntaxError:
        return None
    symbol_map = {s.get("id"): s for s in root.iter(f"{{{SVG_DOC_NS}}}symbol")
                  if s.get("id")}
    _inline_uses(root, symbol_map)

    editable = 0  # editable elements that remain (text + kept painted shapes)
    for el in list(root.iter()):
        ln = _ln(el)
        if any(_ln(a) in _DEF_CONTAINERS for a in el.iterancestors()):
            continue  # template geometry inside <defs>/<symbol>, never painted directly
        if ln == "text":
            editable += 1  # text is always editable; it stays in the layer untouched
            continue
        if ln not in _PAINTED_TAGS:
            continue
        if _is_atmosphere(el):
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)
        else:
            if el.get("filter"):
                del el.attrib["filter"]  # soft shadow stays in background picture A
            editable += 1

    if editable == 0:
        return None  # nothing editable to layer — keep the plain single picture
    return etree.tostring(root, xml_declaration=True, encoding="utf-8")


def add_svg_slide(prs, svg_path: Path, png_path: Path, embed_svg: bool,
                  pic_path: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    pic = slide.shapes.add_picture(
        str(pic_path or png_path), 0, 0,
        width=prs.slide_width, height=prs.slide_height,
    )
    if embed_svg:
        svg_bytes = svg_path.read_bytes()
        svg_rId = _add_svg_part(slide.part, svg_bytes)
        _inject_svg_ext(pic._element, svg_rId)
    return slide


def add_decomposed_slide(prs, full_png: Path, content_png: Path,
                         content_svg_bytes: bytes):
    """Two-picture slide: a full-render background picture (movable/resizable,
    shows in every viewer) plus a transparent content layer carrying the
    editable svgBlip vector on top. See the module note above build_editable_layer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # A — full render, the one movable background image. No svgBlip, so
    #     PowerPoint keeps showing the raster (atmosphere, shadows and all).
    slide.shapes.add_picture(str(full_png), 0, 0,
                             width=prs.slide_width, height=prs.slide_height)
    # B — content-only layer; PowerPoint renders the svgBlip vector over A.
    pic_b = slide.shapes.add_picture(str(content_png), 0, 0,
                                     width=prs.slide_width, height=prs.slide_height)
    svg_rId = _add_svg_part(slide.part, content_svg_bytes)
    _inject_svg_ext(pic_b._element, svg_rId)
    return slide


def strip_notes_infrastructure(pptx_path: Path) -> None:
    """python-pptx's notes_text_frame mechanism leaves Content_Types overrides
    and slide rels that Keynote refuses to parse — it shows "invalid file
    format" and won't open the file at all.

    We strip ALL notes-related parts and references after python-pptx saves,
    leaving a clean image-based PPTX that opens in Keynote / PowerPoint /
    Preview / Quick Look uniformly. Notes content is preserved in a sibling
    `<stem>.notes.md` file by save_notes_sidecar() — not in the .pptx itself.

    Skipped when --keep-notes is set (PowerPoint-only workflows).
    """
    import re
    import tempfile
    import shutil
    import zipfile

    tmp_path = pptx_path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(pptx_path) as zin, \
         zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.namelist():
            # Drop notes-related parts entirely
            if ("notesSlide" in item or
                "notesMaster" in item or
                item.endswith("theme2.xml")):
                continue
            data = zin.read(item)
            if item == "[Content_Types].xml":
                text = data.decode("utf-8")
                text = re.sub(r"<Override[^>]*notesSlide[^>]*/>", "", text)
                text = re.sub(r"<Override[^>]*notesMaster[^>]*/>", "", text)
                # Also drop the theme2 override since we drop the file
                text = re.sub(r'<Override[^>]*PartName="/ppt/theme/theme2\.xml"[^>]*/>',
                              "", text)
                data = text.encode("utf-8")
            if item.endswith(".rels"):
                text = data.decode("utf-8")
                text = re.sub(r"<Relationship[^>]*notesSlide[^>]*/>", "", text)
                text = re.sub(r"<Relationship[^>]*notesMaster[^>]*/>", "", text)
                data = text.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(tmp_path, pptx_path)


def save_notes_sidecar(pptx_path: Path, planning: dict) -> Path | None:
    """Write speaker notes as a sibling `.notes.md` file so the planner's
    speaker_notes aren't lost when we strip notes from the .pptx for Keynote
    compatibility. Returns the sidecar path, or None if there were no notes."""
    pages = planning.get("pages") or []
    notes = [(i + 1, p.get("speaker_notes", "")) for i, p in enumerate(pages)
             if p.get("speaker_notes")]
    if not notes:
        return None
    out = pptx_path.with_suffix(".notes.md")
    lines = [f"# Speaker notes — {pptx_path.stem}", ""]
    for idx, note in notes:
        lines.append(f"## Slide {idx}")
        lines.append("")
        lines.append(note.strip())
        lines.append("")
    out.write_text("\n".join(lines), encoding="utf-8")
    return out


def build(pages: list[Path], out_path: Path, planning: dict,
          workdir: Path, png_width: int,
          placeholder_only: bool, embed_svg: bool,
          also_pdf: bool, pdf_path: Path | None,
          keep_notes: bool,
          no_anim: bool = False, gif_width: int = DEFAULT_GIF_WIDTH,
          no_decompose: bool = False):
    workdir.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU

    placeholder_path = workdir / "_placeholder.png"
    if not placeholder_path.exists():
        placeholder_path.write_bytes(_PLACEHOLDER_PNG_BYTES)

    real_pngs = []   # ordered list of per-slide PNG paths used for the optional PDF
    slide_objs = []  # ordered list of slide objects, for the notes second-pass
    engine_reported = False
    used_placeholder_anywhere = False
    for i, svg_path in enumerate(pages):
        if placeholder_only:
            png_path = placeholder_path
            used_placeholder_anywhere = True
        else:
            png_path = workdir / f"{svg_path.stem}.png"
            engine = render_real_png(svg_path, png_path, width=png_width)
            if engine is None:
                if not engine_reported:
                    print("⚠️  No SVG renderer found. "
                          "Falling back to a 1×1 transparent placeholder PNG.",
                          file=sys.stderr)
                    print("    The .pptx will look BLANK in Keynote / macOS Preview / "
                          "older PowerPoint. Easiest fix:",
                          file=sys.stderr)
                    print("      pip install resvg-py --break-system-packages",
                          file=sys.stderr)
                    print("    (zero system deps; prebuilt wheels for all major platforms.)",
                          file=sys.stderr)
                    engine_reported = True
                png_path = placeholder_path
                used_placeholder_anywhere = True
            elif not engine_reported:
                print(f"PNG fallback engine: {engine} ({png_width}px wide)")
                engine_reported = True

        # Flow-anim pages become looping GIF slides (svgBlip skipped — PowerPoint
        # would otherwise render the static SVG layer on top of the GIF).
        # Only attempted when a real static render succeeded; the static PNG is
        # still what the PDF uses.
        gif_path = None
        if not no_anim and png_path is not placeholder_path:
            if _flow_anim_is_animatable(svg_path.read_text(encoding="utf-8"),
                                        svg_path.name):
                candidate = workdir / f"{svg_path.stem}.gif"
                if render_flow_anim_gif(svg_path, candidate, width=gif_width):
                    gif_path = candidate

        real_pngs.append(png_path)
        if gif_path is not None:
            print(f"[{i+1}/{len(pages)}] {svg_path.name} → slide "
                  "(animated GIF; not Convert-to-Shape editable)", flush=True)
            slide = add_svg_slide(prs, svg_path, png_path, embed_svg=False,
                                  pic_path=gif_path)
        else:
            # Default: split into a movable background picture + an editable
            # content layer (icons inlined, atmosphere/shadows left to the
            # background). Falls back to the plain single-picture slide whenever
            # decomposition isn't applicable or fails, so behavior never regresses.
            content_svg_bytes = None
            content_png = None
            if embed_svg and not no_decompose and png_path is not placeholder_path:
                content_svg_bytes = build_editable_layer(svg_path.read_bytes())
                if content_svg_bytes is not None:
                    content_svg_file = workdir / f"{svg_path.stem}.content.svg"
                    content_svg_file.write_bytes(content_svg_bytes)
                    candidate_png = workdir / f"{svg_path.stem}.content.png"
                    if render_real_png(content_svg_file, candidate_png, width=png_width):
                        content_png = candidate_png
                    else:
                        content_svg_bytes = None  # render failed → plain slide
            if content_svg_bytes is not None and content_png is not None:
                print(f"[{i+1}/{len(pages)}] {svg_path.name} → slide "
                      "(background image + editable content layer)", flush=True)
                slide = add_decomposed_slide(prs, png_path, content_png,
                                             content_svg_bytes)
            else:
                print(f"[{i+1}/{len(pages)}] {svg_path.name} → slide", flush=True)
                slide = add_svg_slide(prs, svg_path, png_path, embed_svg=embed_svg)
        slide_objs.append(slide)

    # Speaker notes — only if --keep-notes is set. By default we strip notes
    # entirely from the .pptx and save them to a sibling .notes.md, because
    # python-pptx's notes infrastructure makes Keynote reject the file.
    if keep_notes:
        for i, slide in enumerate(slide_objs):
            notes = speaker_notes_for(i, planning)
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(out_path))

    if not keep_notes:
        # Always strip after save — python-pptx may have added stray notes
        # infrastructure (notesMaster, theme2) even when no notes were set.
        strip_notes_infrastructure(out_path)

    print(f"✅ Wrote PPTX: {out_path}")
    if embed_svg:
        print("   (In PowerPoint 2016+: right-click the content layer → Convert to "
              "Shape to edit text/cards/icons; the background image moves separately.)")

    # Save sidecar notes file (regardless of --keep-notes, so notes are always
    # recoverable from a human-readable source).
    notes_md = save_notes_sidecar(out_path, planning)
    if notes_md:
        print(f"📝 Wrote notes: {notes_md}")
    if used_placeholder_anywhere and not placeholder_only:
        print("⚠️  Some slides used the 1×1 placeholder because no SVG renderer was "
              "available. Re-run after installing one for proper Keynote / Preview display.")

    # Emit the companion PDF if requested (best-effort — several conditions
    # legitimately skip it). Whether or not the PDF is produced, the
    # deliverables footer below always prints, so the caller reads the exact
    # file set from stdout instead of inferring it from which branch ran.
    pdf_out = None
    if also_pdf:
        if not _HAS_IMG2PDF:
            print("⚠️  Skipping PDF: img2pdf not installed. To enable, run:",
                  file=sys.stderr)
            print("      pip install img2pdf --break-system-packages",
                  file=sys.stderr)
        elif placeholder_only:
            # User explicitly chose 1×1 transparent placeholders — a PDF built
            # from those would be entirely blank, so skip rather than mislead.
            print("⚠️  Skipping PDF: --placeholder-only mode produces no real "
                  "PNG renders, so the PDF would be blank.", file=sys.stderr)
        elif used_placeholder_anywhere:
            # Some slides fell back to placeholder because the renderer failed
            # mid-stream — PDF would be partly blank, skip it.
            print("⚠️  Skipping PDF: no real PNG renders are available "
                  "(no SVG renderer installed).", file=sys.stderr)
        else:
            pdf_out = pdf_path or out_path.with_suffix(".pdf")
            write_pdf(real_pngs, pdf_out)
            print(f"✅ Wrote PDF:  {pdf_out}")

    # Deliverables footer — printed unconditionally so the caller can read the
    # exact file set from stdout regardless of whether the PDF was produced.
    # Only files that actually exist on disk are listed.
    notes_sidecar = out_path.with_suffix(".notes.md")
    deliverables = [out_path]
    if pdf_out is not None and pdf_out.exists():
        deliverables.append(pdf_out)
    if notes_sidecar.exists():
        deliverables.append(notes_sidecar)
    print()
    print("=" * 70)
    print(f"⚠️  CRITICAL: {len(deliverables)} file(s) produced. DELIVER ALL OF THEM.")
    print("   The user expects every one of these in their chat / Downloads:")
    for p in deliverables:
        print(f"     • {p}")
    print("   Delivering only the .pptx and dropping the .pdf/.notes.md")
    print("   is the #1 DeckForge bug. Do not do this.")
    print("=" * 70)


def write_pdf(png_paths: list[Path], pdf_path: Path) -> None:
    """Assemble per-slide PNGs into a single multi-page PDF via img2pdf."""
    pdf_path.parent.mkdir(parents=True, exist_ok=True)
    # img2pdf accepts a list of file paths or bytes; paths is simpler.
    paths_str = [str(p) for p in png_paths]
    # Set page size to match a standard 16:9 slide so the PDF previews at the
    # right aspect on any reader. img2pdf computes from image DPI by default;
    # we override via layout_fun for an exact 13.333" × 7.5" page.
    layout = img2pdf.get_layout_fun(
        pagesize=(img2pdf.in_to_pt(13.333), img2pdf.in_to_pt(7.5)),
    )
    with open(pdf_path, "wb") as f:
        f.write(img2pdf.convert(paths_str, layout_fun=layout))


def main():
    p = argparse.ArgumentParser(description=__doc__,
                                formatter_class=argparse.RawDescriptionHelpFormatter)
    p.add_argument("--pages-dir", required=True, help="Directory containing page_*.svg files")
    p.add_argument("--output", required=True, help="Output .pptx path")
    p.add_argument("--planning", default=None, help="planning.json path (auto-detected if omitted)")
    p.add_argument("--no-svg", action="store_true",
                   help="Skip the svgBlip extension; slides become image-only. "
                        "Use this for viewers that choke on the SVG ext.")
    p.add_argument("--placeholder-only", action="store_true",
                   help="Force the 1×1 transparent placeholder PNG fallback even "
                        "when a real renderer is available. Smaller file but only "
                        "displays correctly in PowerPoint 2016+.")
    p.add_argument("--png-width", type=int, default=2560,
                   help="Width (px) of the PNG fallback. Default 2560 (2× DPI).")
    p.add_argument("--no-anim", action="store_true",
                   help="Disable flow-edge animation. By default, pages marked "
                        "with class=\"flow-anim\" dashed edges are embedded as "
                        "looping animated GIFs (those slides skip svgBlip and "
                        "are not Convert-to-Shape editable).")
    p.add_argument("--gif-width", type=int, default=DEFAULT_GIF_WIDTH,
                   help="Width (px) of animated GIF frames. Default 1600.")
    p.add_argument("--no-decompose", action="store_true",
                   help="Embed each slide as ONE picture (full render + svgBlip) "
                        "instead of the default background-image + editable-content-"
                        "layer split. The split makes icons/cards/lines individually "
                        "movable after Convert to Shape; use this flag to revert to "
                        "the older single-picture behavior.")
    p.add_argument("--no-pdf", action="store_true",
                   help="Skip the companion .pdf (by default, a PDF with the same "
                        "stem as --output is written alongside the .pptx).")
    p.add_argument("--pdf-output", default=None,
                   help="Explicit PDF path. If omitted, derived from --output by "
                        "swapping .pptx → .pdf.")
    p.add_argument("--keep-notes", action="store_true",
                   help="Embed speaker notes inside the .pptx via python-pptx. "
                        "Default is OFF because python-pptx's notes infrastructure "
                        "makes Keynote refuse the file. Notes are always written "
                        "to a sibling <stem>.notes.md regardless of this flag.")
    p.add_argument("--workdir", default=None,
                   help="Directory for intermediate files (default: <pages-dir>/_renders)")
    args = p.parse_args()

    pages_dir = Path(args.pages_dir).resolve()
    out_path = Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    planning_path = Path(args.planning).resolve() if args.planning else None
    workdir = Path(args.workdir).resolve() if args.workdir else pages_dir / "_renders"
    pdf_path = Path(args.pdf_output).resolve() if args.pdf_output else None

    pages = collect_pages(pages_dir)
    planning = load_planning(planning_path, pages_dir)

    # Fail fast: --no-svg requires a working SVG renderer. Without one, every
    # slide would be just the 1×1 transparent placeholder PNG → an all-blank
    # PPTX in every viewer. That's the worst kind of silent failure.
    if args.no_svg and not args.placeholder_only:
        engine = verify_renderer_works(pages[0])
        if engine is None:
            print("❌ --no-svg requires an SVG renderer, but none is available.",
                  file=sys.stderr)
            print("   Without one, every slide will be a 1×1 transparent PNG "
                  "(blank everywhere).", file=sys.stderr)
            print("   Pick one of these fixes:", file=sys.stderr)
            print("     pip install resvg-py --break-system-packages",
                  file=sys.stderr)
            print("     brew install librsvg          (macOS)", file=sys.stderr)
            print("     apt-get install librsvg2-bin  (Linux)", file=sys.stderr)
            print("   Or remove --no-svg — the default svgBlip mode embeds "
                  "the SVG so modern PowerPoint can render it directly.",
                  file=sys.stderr)
            sys.exit(2)

    build(pages, out_path, planning, workdir, args.png_width,
          placeholder_only=args.placeholder_only,
          embed_svg=not args.no_svg,
          also_pdf=not args.no_pdf,
          pdf_path=pdf_path,
          keep_notes=args.keep_notes,
          no_anim=args.no_anim,
          gif_width=args.gif_width,
          no_decompose=args.no_decompose)


if __name__ == "__main__":
    main()
